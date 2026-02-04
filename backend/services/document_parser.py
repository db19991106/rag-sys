from typing import Optional, Dict, Any, Tuple
from pathlib import Path
import pypdf
from docx import Document as DocxDocument
import markdown
import csv
import json
import html2text
import xml.etree.ElementTree as ET
from pptx import Presentation
from openpyxl import load_workbook
from utils.logger import logger
from utils.file_utils import get_file_extension
from utils.retry import retry, DOCUMENT_PARSER_RETRY_CONFIG


class DocumentParser:
    """文档解析器 - 支持多种文档格式，带重试机制"""

    @staticmethod
    @retry(
        max_attempts=DOCUMENT_PARSER_RETRY_CONFIG["max_attempts"],
        base_delay=DOCUMENT_PARSER_RETRY_CONFIG["base_delay"],
        backoff_factor=DOCUMENT_PARSER_RETRY_CONFIG["backoff_factor"],
        max_delay=DOCUMENT_PARSER_RETRY_CONFIG["max_delay"],
        exceptions=(IOError, pypdf.errors.PdfReadError, Exception)
    )
    def parse(file_path: str) -> Optional[str]:
        """
        解析文档内容

        Args:
            file_path: 文件路径

        Returns:
            文档文本内容，解析失败返回 None
        """
        ext = get_file_extension(file_path)

        try:
            if ext == '.txt':
                return DocumentParser._parse_txt(file_path)
            elif ext == '.pdf':
                return DocumentParser._parse_pdf(file_path)
            elif ext == '.docx':
                return DocumentParser._parse_docx(file_path)
            elif ext == '.md':
                return DocumentParser._parse_markdown(file_path)
            elif ext == '.csv':
                return DocumentParser._parse_csv(file_path)
            elif ext == '.json':
                return DocumentParser._parse_json(file_path)
            elif ext == '.html' or ext == '.htm':
                return DocumentParser._parse_html(file_path)
            elif ext == '.xml':
                return DocumentParser._parse_xml(file_path)
            elif ext == '.pptx':
                return DocumentParser._parse_pptx(file_path)
            elif ext == '.xlsx':
                return DocumentParser._parse_xlsx(file_path)
            else:
                logger.error(f"不支持的文件格式: {ext}")
                return None
        except Exception as e:
            logger.error(f"解析文件失败 {file_path}: {str(e)}")
            return None

    @staticmethod
    def parse_with_metadata(file_path: str, additional_metadata: Optional[Dict[str, Any]] = None) -> Tuple[Optional[str], Optional[Dict[str, Any]]]:
        """
        解析文档内容并返回元数据

        Args:
            file_path: 文件路径
            additional_metadata: 额外的元数据

        Returns:
            (文档文本内容, 元数据)，解析失败返回 (None, None)
        """
        content = DocumentParser.parse(file_path)
        if not content:
            return None, None

        # 提取元数据
        metadata = {
            'file_path': file_path,
            'file_size': Path(file_path).stat().st_size if Path(file_path).exists() else 0,
            'content_length': len(content) if content else 0,
            'format': get_file_extension(file_path),
            'parsed_at': Path(file_path).stat().st_mtime if Path(file_path).exists() else None
        }

        # 添加额外的元数据
        if additional_metadata:
            metadata.update(additional_metadata)

        return content, metadata

    @staticmethod
    def _parse_txt(file_path: str) -> str:
        """解析 TXT 文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        """解析 PDF 文件"""
        content = []
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
        return '\n'.join(content)

    @staticmethod
    def _parse_docx(file_path: str) -> str:
        """解析 DOCX 文件"""
        doc = DocxDocument(file_path)
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                content.append(paragraph.text)
        return '\n'.join(content)

    @staticmethod
    def _parse_markdown(file_path: str) -> str:
        """解析 Markdown 文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        # 返回原始 Markdown 内容
        return md_content

    @staticmethod
    def _parse_csv(file_path: str) -> str:
        """解析 CSV 文件"""
        content = []
        with open(file_path, 'r', encoding='utf-8', newline='') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                if row:
                    content.append(','.join(row))
        return '\n'.join(content)

    @staticmethod
    def _parse_json(file_path: str) -> str:
        """解析 JSON 文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        # 将 JSON 转换为格式化字符串
        return json.dumps(data, ensure_ascii=False, indent=2)

    @staticmethod
    def _parse_html(file_path: str) -> str:
        """解析 HTML 文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        # 将 HTML 转换为纯文本
        h = html2text.HTML2Text()
        h.ignore_links = True
        h.ignore_images = True
        return h.handle(html_content)

    @staticmethod
    def _parse_xml(file_path: str) -> str:
        """解析 XML 文件"""
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        def element_to_text(element, indent=0):
            text = []
            indent_str = '  ' * indent
            text.append(f"{indent_str}{element.tag}:")
            
            if element.text and element.text.strip():
                text.append(f"{indent_str}  {element.text.strip()}")
            
            for child in element:
                text.extend(element_to_text(child, indent + 1))
            
            if element.tail and element.tail.strip():
                text.append(f"{indent_str}  {element.tail.strip()}")
            
            return text
        
        return '\n'.join(element_to_text(root))

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        """解析 PPTX 文件"""
        prs = Presentation(file_path)
        content = []
        
        for i, slide in enumerate(prs.slides, 1):
            content.append(f"幻灯片 {i}:")
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text:
                        content.append(f"  {shape.text}")
            content.append('')
        
        return '\n'.join(content)

    @staticmethod
    def _parse_xlsx(file_path: str) -> str:
        """解析 XLSX 文件"""
        wb = load_workbook(file_path)
        content = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            content.append(f"工作表: {sheet_name}")
            
            # 读取前100行数据
            for row in ws.iter_rows(min_row=1, max_row=100, values_only=True):
                if any(cell is not None for cell in row):
                    row_content = [str(cell) if cell is not None else '' for cell in row]
                    content.append('  ' + '\t'.join(row_content))
            
            content.append('')
        
        return '\n'.join(content)